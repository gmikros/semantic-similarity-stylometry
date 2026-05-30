from __future__ import annotations

from pathlib import Path
import re

from docx import Document
from pptx import Presentation


ROOT = Path(r"C:\Users\USER01\Dropbox\Workplace\D\George\PAPERS\Qualico 2025\Mikros, Cech")
PPTX_PATH = ROOT / "Investigating the Impact of Semantic Similarity on Stylometric Attribution Using Controlled Artificial Texts and Delta Distances.pptx"
PAPER_PATH = ROOT / "Paper" / "260527_mikros_etal_semantic_smililarity.docx"
OUT = ROOT / "analysis" / "output" / "doc" / "rewrite_sources_summary.md"


def norm(s: str) -> str:
    return re.sub(r"\s+", " ", s or "").strip()


def extract_pptx() -> list[str]:
    prs = Presentation(PPTX_PATH)
    lines: list[str] = ["## Presentation Notes", ""]
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                txt = norm(shape.text)
                if not txt:
                    continue
                if not title:
                    title = txt
                else:
                    bullets.append(txt)
        lines.append(f"### Slide {i}: {title or '[No title]'}")
        for b in bullets[:8]:
            lines.append(f"- {b}")
        lines.append("")
    return lines


def extract_docx() -> list[str]:
    doc = Document(PAPER_PATH)
    lines: list[str] = ["## Paper Draft Notes", ""]
    sections = {"Introduction": [], "Methods": [], "Results": []}
    current = None
    for p in doc.paragraphs:
        t = norm(p.text)
        if not t:
            continue
        lower = t.lower()
        if "introduction" == lower or lower.startswith("1. introduction"):
            current = "Introduction"
        elif "methods" == lower or lower.startswith("2. methods"):
            current = "Methods"
        elif "results" == lower or lower.startswith("3. results"):
            current = "Results"
        if current in sections:
            sections[current].append(t)

    for name, paras in sections.items():
        lines.append(f"### {name}")
        for p in paras[:30]:
            lines.append(f"- {p}")
        lines.append("")
    return lines


def main() -> None:
    OUT.parent.mkdir(parents=True, exist_ok=True)
    content = []
    content.extend(extract_pptx())
    content.extend(extract_docx())
    OUT.write_text("\n".join(content), encoding="utf-8")
    print(str(OUT))


if __name__ == "__main__":
    main()
